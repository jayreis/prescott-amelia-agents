#!/usr/bin/env python3
"""
Prospect Audit Deck — personal brand template.

Usage: copy this file to {brand_slug}_audit_deck.py, edit CONFIG below
using findings from prospect_audit_checklist.md, then run it.

The CONFIG dict below is filled with example placeholder data so the
script runs out of the box — replace every value before sending to a
real prospect.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

# ── Brand colors — replace with your own ──────────────────────
NAVY       = RGBColor(0x05, 0x09, 0x1A)   # deep navy background
NAVY_MID   = RGBColor(0x0B, 0x13, 0x26)
NAVY_LIGHT = RGBColor(0x10, 0x1D, 0x38)
GOLD       = RGBColor(0xC9, 0xA8, 0x5C)
GOLD_DARK  = RGBColor(0x9B, 0x7F, 0x40)
CREAM      = RGBColor(0xF6, 0xF2, 0xE9)
CREAM_DARK = RGBColor(0xEB, 0xE5, 0xD6)
BODY_LIGHT = RGBColor(0xD9, 0xD4, 0xC9)   # body text on dark bg
MUTED      = RGBColor(0x7B, 0x84, 0x9A)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
DGRAY      = RGBColor(0x33, 0x33, 0x33)
RED        = RGBColor(0xB5, 0x4A, 0x3C)
GREEN      = RGBColor(0x5C, 0x8A, 0x5E)
YELLOW     = GOLD

FONT_HEAD = "Cormorant Garamond"
FONT_BODY = "Outfit"

# ═══════════════════════════════════════════════════════════
# CONFIG — replace with real audit findings per brand
# ═══════════════════════════════════════════════════════════
CONFIG = {
    "brand_name": "Example Brand Co.",
    "prepared_for": "Director of Operations",
    "audit_date": "July 2026",
    "output_path": "./EXAMPLE_prospect_audit_deck.pptx",

    "scorecard": [
        ("2", "of 7 Core Flows Detected", RED),
        ("Confirmation only", "Welcome Series Depth", RED),
        ("48hr", "Time to First Cart Email", GOLD_DARK),
        ("3", "Campaign Emails in 14 Days", GREEN),
    ],

    "key_findings": [
        ("GAP", "No welcome series beyond the signup confirmation — the highest-ROI automation in ecommerce is going unused."),
        ("GAP", "Browse abandonment is not set up — no email received after visiting product pages without adding to cart."),
        ("RISK", "Cart abandon discount escalates from 10% to 20% across the sequence, which tends to train customers to wait for the bigger offer instead of buying at full price."),
        ("OPPORTUNITY", "Campaign cadence is consistent (3 sends/14 days) and reasonably varied — a solid foundation to build flows on top of."),
    ],

    "flow_gaps": [
        ("Welcome Series",        "Partial",      YELLOW, "1 confirmation email only. No brand story, education, or purchase-incentive follow-up detected over 10 days."),
        ("Abandoned Browse",      "Not Detected", RED,    "No email after browsing 3 product pages and leaving without adding to cart."),
        ("Abandoned Cart",        "Detected",     GREEN,  "3-email sequence over 5 days. Discount escalates 10% to 20% — worth revisiting."),
        ("Abandoned Checkout",    "Not Distinct", RED,    "Same sequence as cart abandon — no dedicated checkout-stage messaging."),
        ("Post-Purchase / Review","Not Detected", RED,    "No thank-you or review-request email after a completed order."),
        ("Replenishment/Reorder", "Unknown",      MUTED,  "Consumable product category — worth confirming if a reorder flow exists (needs longer observation window)."),
        ("Winback",               "Unknown",      MUTED,  "Requires 60-90+ days of inactivity to observe — flagged as unconfirmed, not absent."),
    ],

    "welcome_deep_dive": {
        "observed": "Single confirmation email sent immediately on signup. No follow-up over the next 10 days. No brand introduction, no product education, no first-purchase incentive.",
        "benchmark_stats": [
            ("51%", "Average open rate on welcome email #1 (industry benchmark)"),
            ("3x",  "Higher revenue per email vs. standard campaigns"),
            ("3-4", "Emails recommended over 5-10 days per Lifecycle Revenue Playbook"),
        ],
        "recommended_steps": [
            ("Email 1", "Immediately", "Brand welcome + hero product"),
            ("Email 2", "Day 2",       "Product education + social proof"),
            ("Email 3", "Day 5",       "First purchase incentive"),
            ("Email 4", "Day 10",      "Category explore / bestsellers"),
        ],
    },

    "opportunity_estimate": [
        ("Illustrative", "Adding a 4-email welcome series against your current list size, at DTC benchmark rates, is typically one of the highest-return single builds available — actual number depends on list size and AOV, which I don't have visibility into from the outside."),
        ("Illustrative", "Fixing the cart discount escalation (flat offer instead of 10%->20%) typically protects margin without a meaningful drop in recovered revenue."),
        ("Illustrative", "Browse abandonment is usually a smaller but pure-incremental line — it's reaching browsers no other flow currently touches."),
    ],

    "engagement_options": [
        ("Fixed-fee sprint", "2-week segmentation + flow rebuild covering welcome and cart discount fix. Fastest path to seeing results before any ongoing commitment."),
        ("Monthly retainer", "Ongoing flow + campaign management once the foundation is fixed — scoped after the sprint based on what's working."),
    ],

    "contact": {
        "name": "Your Name",
        "title": "Lifecycle Marketing Consultant — CPG & DTC",
        "site": "yoursite.com",
    },
}


def build(cfg):
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    BLANK = prs.slide_layouts[6]

    def bg(slide, color):
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = color

    def box(slide, l, t, w, h, fill_color=None, border_color=None, border_width=Pt(0)):
        shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
        shape.line.fill.background()
        if fill_color:
            shape.fill.solid()
            shape.fill.fore_color.rgb = fill_color
        else:
            shape.fill.background()
        if border_color:
            shape.line.color.rgb = border_color
            shape.line.width = border_width
        else:
            shape.line.fill.background()
        return shape

    def txt(slide, text, l, t, w, h, size=18, bold=False, color=WHITE,
            align=PP_ALIGN.LEFT, wrap=True, font=FONT_BODY, italic=False):
        tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
        tf = tb.text_frame
        tf.word_wrap = wrap
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color
        run.font.name = font
        return tb

    # ── SLIDE 1 — TITLE ──────────────────────────────────────
    s = prs.slides.add_slide(BLANK)
    bg(s, NAVY)
    box(s, 0, 0, 13.33, 0.06, fill_color=GOLD)
    box(s, 0, 7.44, 13.33, 0.06, fill_color=GOLD)
    box(s, 0, 2.5, 13.33, 2.5, fill_color=NAVY_MID)

    txt(s, cfg["brand_name"].upper(), 0.8, 1.15, 11, 0.6, size=13, bold=True,
        color=GOLD, font=FONT_BODY)
    txt(s, "Email & Retention Audit", 0.8, 1.7, 11, 1.0, size=44, bold=True,
        color=WHITE, font=FONT_HEAD)
    txt(s, "A complimentary lifecycle marketing review — flows, cadence, and quick opportunities",
        0.8, 2.75, 11, 0.5, size=15, color=MUTED, font=FONT_BODY)
    txt(s, f"Prepared by {cfg['contact']['name']}  |  {cfg['audit_date']}", 0.8, 6.7, 8, 0.4,
        size=11, color=GOLD_DARK, font=FONT_BODY)

    # ── SLIDE 2 — EXECUTIVE SUMMARY ──────────────────────────
    s = prs.slides.add_slide(BLANK)
    bg(s, CREAM)
    box(s, 0, 0, 13.33, 1.1, fill_color=NAVY)
    txt(s, "EXECUTIVE SUMMARY", 0.5, 0.15, 12, 0.35, size=11, bold=True, color=GOLD)
    txt(s, f"The State of {cfg['brand_name']}'s Email Program", 0.5, 0.5, 12, 0.55,
        size=26, bold=True, color=WHITE, font=FONT_HEAD)

    for i, (val, lbl, color) in enumerate(cfg["scorecard"]):
        x = 0.4 + i * 3.15
        box(s, x, 1.3, 2.9, 1.6, fill_color=color)
        txt(s, str(val), x, 1.45, 2.9, 0.75, size=34, bold=True, color=WHITE,
            align=PP_ALIGN.CENTER, font=FONT_HEAD)
        txt(s, lbl, x, 2.15, 2.9, 0.6, size=10, color=WHITE, align=PP_ALIGN.CENTER)

    txt(s, "KEY FINDINGS", 0.4, 3.1, 12, 0.3, size=9, bold=True, color=DGRAY)
    for i, (tag, finding) in enumerate(cfg["key_findings"]):
        y = 3.45 + i * 0.82
        box(s, 0.4, y, 12.5, 0.72, fill_color=WHITE, border_color=CREAM_DARK, border_width=Pt(1))
        txt(s, tag, 0.6, y + 0.08, 2.0, 0.35, size=9, bold=True, color=NAVY)
        txt(s, finding, 2.6, y + 0.06, 10.1, 0.58, size=11, color=DGRAY)

    # ── SLIDE 3 — FLOW GAP ANALYSIS ──────────────────────────
    s = prs.slides.add_slide(BLANK)
    bg(s, WHITE)
    box(s, 0, 0, 13.33, 1.1, fill_color=NAVY)
    txt(s, "FLOW GAP ANALYSIS", 0.5, 0.15, 12, 0.35, size=11, bold=True, color=GOLD)
    txt(s, "Core Lifecycle Flows — What's There, What's Missing", 0.5, 0.5, 12, 0.55,
        size=24, bold=True, color=WHITE, font=FONT_HEAD)

    for i, (flow, status, color, note) in enumerate(cfg["flow_gaps"]):
        col = i // 4
        row = i % 4
        x = 0.35 + col * 6.45
        y = 1.25 + row * 1.35
        box(s, x, y, 6.2, 1.2, fill_color=CREAM)
        box(s, x, y, 1.7, 1.2, fill_color=color)
        txt(s, status, x, y + 0.4, 1.7, 0.45, size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        txt(s, flow, x + 1.85, y + 0.1, 4.2, 0.4, size=13, bold=True, color=NAVY, font=FONT_HEAD)
        txt(s, note, x + 1.85, y + 0.5, 4.2, 0.65, size=9, color=DGRAY)

    # ── SLIDE 4 — WELCOME FLOW DEEP DIVE ─────────────────────
    wd = cfg["welcome_deep_dive"]
    s = prs.slides.add_slide(BLANK)
    bg(s, CREAM)
    box(s, 0, 0, 13.33, 1.1, fill_color=GOLD_DARK)
    txt(s, "DEEP DIVE — WELCOME SERIES", 0.5, 0.1, 12, 0.35, size=11, bold=True, color=NAVY)
    txt(s, "Highest-ROI Flow, Currently Underbuilt", 0.5, 0.45, 12.5, 0.6, size=24, bold=True,
        color=WHITE, font=FONT_HEAD)

    box(s, 0.4, 1.2, 6.0, 5.9, fill_color=WHITE)
    txt(s, "WHAT I OBSERVED", 0.7, 1.35, 5.5, 0.35, size=10, bold=True, color=GOLD_DARK)
    txt(s, wd["observed"], 0.7, 1.75, 5.5, 1.6, size=12, color=DGRAY)

    box(s, 6.8, 1.2, 6.1, 2.7, fill_color=NAVY)
    txt(s, "INDUSTRY BENCHMARK", 7.0, 1.35, 5.7, 0.35, size=10, bold=True, color=GOLD)
    for i, (v, l) in enumerate(wd["benchmark_stats"]):
        txt(s, v, 7.0, 1.8 + i * 0.75, 2.0, 0.5, size=30, bold=True, color=GOLD, font=FONT_HEAD)
        txt(s, l, 9.2, 1.9 + i * 0.75, 3.5, 0.55, size=10, color=BODY_LIGHT)

    box(s, 6.8, 4.05, 6.1, 3.05, fill_color=WHITE)
    txt(s, "RECOMMENDED SEQUENCE", 7.0, 4.2, 5.7, 0.35, size=10, bold=True, color=GOLD_DARK)
    for i, (e, t, d) in enumerate(wd["recommended_steps"]):
        y = 4.6 + i * 0.58
        box(s, 7.0, y, 0.7, 0.42, fill_color=NAVY)
        txt(s, e, 7.0, y + 0.05, 0.7, 0.32, size=8, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        txt(s, t, 7.8, y + 0.05, 1.3, 0.32, size=9, color=GOLD_DARK)
        txt(s, d, 9.2, y + 0.05, 3.5, 0.32, size=9, color=DGRAY)

    # ── SLIDE 5 — OPPORTUNITY SNAPSHOT ───────────────────────
    s = prs.slides.add_slide(BLANK)
    bg(s, NAVY)
    box(s, 0, 0, 13.33, 1.1, fill_color=NAVY_MID)
    txt(s, "OPPORTUNITY SNAPSHOT", 0.5, 0.15, 12, 0.35, size=11, bold=True, color=GOLD)
    txt(s, "Where the Upside Is — Illustrative, Not a Guarantee", 0.5, 0.5, 12, 0.55,
        size=24, bold=True, color=WHITE, font=FONT_HEAD)

    for i, (tag, body) in enumerate(cfg["opportunity_estimate"]):
        y = 1.35 + i * 1.8
        box(s, 0.4, y, 12.5, 1.6, fill_color=NAVY_LIGHT)
        box(s, 0.4, y, 0.08, 1.6, fill_color=GOLD)
        txt(s, tag.upper(), 0.65, y + 0.15, 3, 0.3, size=9, bold=True, color=GOLD)
        txt(s, body, 0.65, y + 0.5, 11.8, 1.0, size=12, color=BODY_LIGHT)

    # ── SLIDE 6 — HOW I'D APPROACH THIS ──────────────────────
    s = prs.slides.add_slide(BLANK)
    bg(s, CREAM)
    box(s, 0, 0, 13.33, 1.1, fill_color=NAVY)
    txt(s, "HOW I'D APPROACH THIS", 0.5, 0.15, 12, 0.35, size=11, bold=True, color=GOLD)
    txt(s, "Two Ways to Work Together", 0.5, 0.5, 12, 0.55, size=26, bold=True,
        color=WHITE, font=FONT_HEAD)

    for i, (label, body) in enumerate(cfg["engagement_options"]):
        x = 0.5 + i * 6.3
        box(s, x, 1.3, 6.0, 5.6, fill_color=WHITE)
        box(s, x, 1.3, 6.0, 0.7, fill_color=GOLD_DARK)
        txt(s, label, x, 1.42, 6.0, 0.5, size=18, bold=True, color=WHITE,
            align=PP_ALIGN.CENTER, font=FONT_HEAD)
        txt(s, body, x + 0.4, 2.3, 5.2, 4.2, size=13, color=DGRAY)

    # ── SLIDE 7 — NEXT STEPS ─────────────────────────────────
    s = prs.slides.add_slide(BLANK)
    bg(s, NAVY)
    box(s, 0, 0, 13.33, 0.06, fill_color=GOLD)
    box(s, 0, 7.44, 13.33, 0.06, fill_color=GOLD)
    txt(s, "NEXT STEPS", 0.6, 1.1, 12, 0.45, size=13, bold=True, color=GOLD)
    txt(s, "Let's Talk It Through", 0.6, 1.55, 10, 0.8, size=40, bold=True, color=WHITE, font=FONT_HEAD)
    txt(s, "No obligation — this is the audit, not a pitch. If it's useful, we can talk about what's next.",
        0.6, 2.4, 10.5, 0.5, size=14, color=BODY_LIGHT)

    c = cfg["contact"]
    txt(s, c["name"], 0.6, 5.6, 8, 0.5, size=22, bold=True, color=GOLD, font=FONT_HEAD, italic=True)
    txt(s, c["title"], 0.6, 6.1, 8, 0.4, size=13, color=BODY_LIGHT)
    txt(s, c["site"], 0.6, 6.5, 8, 0.4, size=13, color=MUTED)

    # ── Save ──────────────────────────────────────────────────
    out = os.path.expanduser(cfg["output_path"])
    prs.save(out)
    print(f"\nDeck saved: {out}")
    print(f"Slides: {len(prs.slides._sldIdLst)}")
    print(f"Open with: open '{out}'")


if __name__ == "__main__":
    build(CONFIG)
